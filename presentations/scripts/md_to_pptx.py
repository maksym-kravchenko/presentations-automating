#!/usr/bin/env python3
"""
md_to_pptx.py
Converts a Slidev Markdown file into a fully editable PPTX.
All text, bullets, headings and code are real PowerPoint objects.

Usage: python scripts/md_to_pptx.py slides/my-talk.md output/my-talk.pptx
"""

import sys, re
from pathlib import Path

try:
    import yaml
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
except ImportError as e:
    print(f"❌  Missing package: {e}")
    print("    Run: pip install python-pptx pyyaml")
    sys.exit(1)


# ── Color palette (matches style.css) ─────────────────────────────────────
C_PRIMARY   = RGBColor(0x4F, 0x46, 0xE5)  # #4F46E5
C_SECONDARY = RGBColor(0x81, 0x8C, 0xF8)  # #818CF8
C_BG        = RGBColor(0xF7, 0xF8, 0xFC)  # #F7F8FC
C_TEXT      = RGBColor(0x1E, 0x1B, 0x4B)  # #1E1B4B
C_MUTED     = RGBColor(0xC7, 0xD2, 0xFE)  # #C7D2FE
C_ACCENT    = RGBColor(0xE0, 0xE7, 0xFF)  # #E0E7FF
C_SURFACE   = RGBColor(0xFF, 0xFF, 0xFF)  # #FFFFFF
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_BG   = RGBColor(0x37, 0x30, 0xA3)  # dark indigo for cover accent

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)
HEADER_H = Inches(0.48)
FOOTER_H = Inches(0.38)
MARGIN   = Inches(0.65)
GAP      = Inches(0.1)


# ── Helpers ────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color: RGBColor, border_color=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0  # sharp corners
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text,
                size=15, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, wrap=True, font='Segoe UI'):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color or C_TEXT
    return txBox


def add_header(slide, subject, author):
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, C_PRIMARY)
    add_textbox(slide, MARGIN, Inches(0.1), Inches(8), HEADER_H,
                subject, size=10, color=C_WHITE)
    add_textbox(slide, Inches(9.5), Inches(0.1), Inches(3.2), HEADER_H,
                author, size=10, color=C_WHITE, align=PP_ALIGN.RIGHT)


def add_footer(slide, company, date, page_num, total):
    y = SLIDE_H - FOOTER_H
    add_rect(slide, 0, y, SLIDE_W, FOOTER_H, C_ACCENT)
    add_textbox(slide, MARGIN, y + Inches(0.07), Inches(4.5), FOOTER_H,
                company, size=9, color=C_PRIMARY)
    add_textbox(slide, Inches(5.5), y + Inches(0.07), Inches(2.5), FOOTER_H,
                date, size=9, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.5), y + Inches(0.07), Inches(2.2), FOOTER_H,
                f"{page_num} / {total}", size=9, color=C_PRIMARY, align=PP_ALIGN.RIGHT)


# ── Markdown parsing ───────────────────────────────────────────────────────

def parse_frontmatter(text):
    text = text.replace('\r\n', '\n')
    match = re.match(r'^---\n(.*?)\n---\n?(.*)', text, re.DOTALL)
    if match:
        try:
            fm = yaml.safe_load(match.group(1)) or {}
        except Exception:
            fm = {}
        return fm, match.group(2).strip()
    return {}, text.strip()


def split_slides(content):
    parts = re.split(r'\n---\n', '\n' + content)
    return [p.strip() for p in parts if p.strip()]


def parse_slide_meta(text):
    text = text.strip()
    match = re.match(r'^---\n(.*?)\n---\n?(.*)', text, re.DOTALL)
    if match:
        try:
            meta = yaml.safe_load(match.group(1)) or {}
        except Exception:
            meta = {}
        return meta, match.group(2).strip()
    return {}, text


def clean_md(text):
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*',     r'\1', text)
    text = re.sub(r'`(.*?)`',       r'\1', text)
    text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
    text = re.sub(r'~~(.*?)~~',     r'\1', text)
    text = re.sub(r'\$.*?\$',       '',    text)
    return text.strip()


def parse_content(text):
    # Strip comments, Vue/Slidev tags, style/script blocks
    text = re.sub(r'<!--.*?-->', '', text, flags=re.DOTALL)
    text = re.sub(r'</?v-clicks?[^>]*>', '', text)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
    text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL)
    text = re.sub(r'<[^>]+>', '', text)

    # Two-column layout
    if '::right::' in text:
        left_raw  = text.split('::right::')[0].replace('::left::', '').strip()
        right_raw = text.split('::right::')[1].strip()
        return [{'type': 'two_cols',
                 'left':  parse_content(left_raw),
                 'right': parse_content(right_raw)}]

    elements = []
    lines = text.split('\n')
    i = 0
    in_code = False
    code_lines = []
    code_lang = ''

    while i < len(lines):
        line = lines[i]

        if re.match(r'^```', line):
            if not in_code:
                in_code   = True
                code_lang = re.sub(r'^```', '', line).split('{')[0].strip()
                code_lines = []
            else:
                in_code = False
                elements.append({'type': 'code', 'lang': code_lang,
                                  'text': '\n'.join(code_lines)})
            i += 1
            continue

        if in_code:
            code_lines.append(line)
            i += 1
            continue

        if line.strip() in ('::left::', '::right::'):
            i += 1
            continue

        if   line.startswith('### '):
            elements.append({'type': 'h3',        'text': line[4:].strip()})
        elif line.startswith('## '):
            elements.append({'type': 'h2',        'text': line[3:].strip()})
        elif line.startswith('# '):
            elements.append({'type': 'h1',        'text': line[2:].strip()})
        elif re.match(r'^  +[-*]\s+', line):
            elements.append({'type': 'bullet',    'text': re.sub(r'^  +[-*]\s+', '', line), 'level': 1})
        elif re.match(r'^[-*]\s+', line):
            elements.append({'type': 'bullet',    'text': re.sub(r'^[-*]\s+', '', line),    'level': 0})
        elif re.match(r'^\d+\.\s+', line):
            elements.append({'type': 'bullet',    'text': re.sub(r'^\d+\.\s+', '', line),   'level': 0})
        elif line.startswith('> '):
            elements.append({'type': 'quote',     'text': line[2:].strip()})
        elif line.strip():
            elements.append({'type': 'paragraph', 'text': line.strip()})

        i += 1

    return elements


# ── Content renderer ───────────────────────────────────────────────────────

def render_elements(slide, elements, x, y, w, h):
    cy = y

    for el in elements:
        if cy >= y + h - Inches(0.3):
            break

        t    = el.get('type')
        text = clean_md(el.get('text', ''))

        if t == 'h1':
            add_textbox(slide, x, cy, w, Inches(0.75), text,
                        size=34, bold=True, color=C_PRIMARY)
            cy += Inches(0.82) + GAP

        elif t == 'h2':
            add_textbox(slide, x, cy, w, Inches(0.58), text,
                        size=24, bold=True, color=C_SECONDARY)
            cy += Inches(0.65) + GAP

        elif t == 'h3':
            add_textbox(slide, x, cy, w, Inches(0.44), text,
                        size=18, bold=True, color=C_PRIMARY)
            cy += Inches(0.50) + GAP

        elif t == 'bullet':
            indent = Inches(0.3) * el.get('level', 0)
            bx = x + Inches(0.15) + indent
            bw = w - Inches(0.15) - indent

            box = slide.shapes.add_textbox(bx, cy, bw, Inches(0.38))
            tf  = box.text_frame
            tf.word_wrap = True
            p   = tf.paragraphs[0]

            r1 = p.add_run()
            r1.text           = '▸  '
            r1.font.size      = Pt(13)
            r1.font.bold      = True
            r1.font.color.rgb = C_SECONDARY
            r1.font.name      = 'Segoe UI'

            r2 = p.add_run()
            r2.text           = text
            r2.font.size      = Pt(15)
            r2.font.color.rgb = C_TEXT
            r2.font.name      = 'Segoe UI'

            cy += Inches(0.40) + GAP

        elif t == 'code':
            code_text  = el.get('text', '')
            code_lines = code_text.split('\n')
            # height based on line count
            code_h = Inches(0.26) * max(len(code_lines), 1) + Inches(0.28)
            code_h = min(code_h, h - (cy - y) - Inches(0.1))  # clamp to remaining space

            add_rect(slide, x, cy, w, code_h, C_SURFACE, border_color=C_MUTED)

            box = slide.shapes.add_textbox(
                x + Inches(0.2), cy + Inches(0.12),
                w - Inches(0.4), code_h - Inches(0.22))
            tf = box.text_frame
            tf.word_wrap = False
            for j, cl in enumerate(code_lines):
                p   = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text           = cl
                run.font.size      = Pt(11)
                run.font.name      = 'Consolas'
                run.font.color.rgb = C_PRIMARY

            cy += code_h + GAP

        elif t == 'quote':
            add_rect(slide, x, cy, Inches(0.07), Inches(0.45), C_SECONDARY)
            add_rect(slide, x, cy, w, Inches(0.45), C_ACCENT)
            add_textbox(slide, x + Inches(0.22), cy + Inches(0.07),
                        w - Inches(0.28), Inches(0.36),
                        text, size=14, italic=True, color=C_PRIMARY)
            cy += Inches(0.52) + GAP

        elif t == 'paragraph':
            add_textbox(slide, x, cy, w, Inches(0.42), text,
                        size=15, color=C_TEXT, wrap=True)
            cy += Inches(0.46) + GAP

        elif t == 'two_cols':
            col_w = (w - Inches(0.4)) / 2
            render_elements(slide, el['left'],
                            x, y, col_w, h)
            render_elements(slide, el['right'],
                            x + col_w + Inches(0.4), y, col_w, h)
            break


# ── Main builder ───────────────────────────────────────────────────────────

def build_pptx(md_path: str, out_path: str):
    raw = Path(md_path).read_text(encoding='utf-8')
    global_fm, body = parse_frontmatter(raw)

    author  = global_fm.get('author',  'Author')
    date    = global_fm.get('date',    '')
    subject = global_fm.get('subject', global_fm.get('title', 'Presentation'))
    company = global_fm.get('company', '')

    slide_texts = split_slides(body)
    total       = len(slide_texts)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # blank layout

    for idx, slide_text in enumerate(slide_texts):
        page_num    = idx + 1
        meta, sbody = parse_slide_meta(slide_text)
        layout      = meta.get('layout', 'default')
        elements    = parse_content(sbody)
        slide       = prs.slides.add_slide(blank)

        # ── Cover ──────────────────────────────────────────────────────────
        if layout == 'cover':
            set_bg(slide, C_PRIMARY)
            add_rect(slide, 0, SLIDE_H - Inches(1.8), SLIDE_W, Inches(1.8), C_DARK_BG)

            title = next((clean_md(e['text']) for e in elements if e['type'] == 'h1'), subject)

            box = slide.shapes.add_textbox(MARGIN, Inches(1.6), SLIDE_W - MARGIN * 2, Inches(2.4))
            tf  = box.text_frame
            tf.word_wrap = True
            p   = tf.paragraphs[0]
            run = p.add_run()
            run.text           = title
            run.font.size      = Pt(46)
            run.font.bold      = True
            run.font.color.rgb = C_WHITE
            run.font.name      = 'Segoe UI'

            # Divider
            add_rect(slide, MARGIN, Inches(4.3), Inches(1.8), Inches(0.04), C_MUTED)

            subtitle = f"{author}  ·  {company}" if company else author
            add_textbox(slide, MARGIN, Inches(4.55), SLIDE_W - MARGIN * 2, Inches(0.55),
                        subtitle, size=17, color=C_MUTED)

            add_textbox(slide, SLIDE_W - Inches(3.5), Inches(5.4),
                        Inches(3.2), Inches(0.4), date,
                        size=13, color=C_MUTED, align=PP_ALIGN.RIGHT)

        # ── Center ─────────────────────────────────────────────────────────
        elif layout == 'center':
            set_bg(slide, C_BG)
            add_header(slide, subject, author)
            add_footer(slide, company, date, page_num, total)
            render_elements(slide, elements,
                            MARGIN,
                            HEADER_H + Inches(1.4),
                            SLIDE_W - MARGIN * 2,
                            SLIDE_H - HEADER_H - FOOTER_H - Inches(2.0))

        # ── Default ────────────────────────────────────────────────────────
        else:
            set_bg(slide, C_BG)
            add_header(slide, subject, author)
            add_footer(slide, company, date, page_num, total)
            render_elements(slide, elements,
                            MARGIN,
                            HEADER_H + Inches(0.28),
                            SLIDE_W - MARGIN * 2,
                            SLIDE_H - HEADER_H - FOOTER_H - Inches(0.38))

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"✅  Saved: {out_path}  ({total} slides)")


# ── Entry point ────────────────────────────────────────────────────────────

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python scripts/md_to_pptx.py slides/my-talk.md output/my-talk.pptx")
        sys.exit(1)
    build_pptx(sys.argv[1], sys.argv[2])